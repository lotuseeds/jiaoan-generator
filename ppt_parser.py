"""
PPT 解析模块
提取PPT中的文字内容和关键图片
"""
import io
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


def _render_pptx_to_screenshots(pptx_path: str, output_dir: str) -> list:
    """
    用 win32com 将 PPTX 导出为临时 PDF，再用 fitz 逐页截图。
    返回 [(slide_index, image_path), ...] 列表（1-based）。
    失败时返回空列表，调用方自动降级为提取内嵌图片。
    """
    import tempfile
    try:
        import fitz
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        pptx_abs = os.path.abspath(pptx_path)
        pdf_path = os.path.join(tempfile.gettempdir(), "jiaoan_slides_tmp.pdf")

        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = False
        prs = app.Presentations.Open(pptx_abs, WithWindow=False)
        prs.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
        prs.Close()
        app.Quit()

        doc = fitz.open(pdf_path)
        results = []
        for i, page in enumerate(doc):
            img_path = os.path.join(output_dir, f"slide_{i+1:02d}.png")
            pix = page.get_pixmap(dpi=150)
            pix.save(img_path)
            results.append((i + 1, img_path))
        doc.close()
        try:
            os.remove(pdf_path)
        except Exception:
            pass
        return results
    except Exception:
        return []


def parse_ppt(ppt_path: str) -> dict:
    """
    解析PPT文件，返回结构化内容
    Returns:
        {
            "slides": [
                {
                    "index": 1,
                    "title": "...",
                    "text": "...",
                    "has_image": True/False,
                    "image_path": "path/to/img.png" or None
                },
                ...
            ],
            "full_text": "所有文字拼接",
            "image_paths": ["path1", "path2", ...]
        }
    """
    prs = Presentation(ppt_path)
    output_dir = os.path.join(os.path.dirname(ppt_path), "_ppt_images")
    os.makedirs(output_dir, exist_ok=True)

    slides_data = []
    all_text_parts = []
    image_paths = []

    # 优先尝试全页截图（win32com + fitz），失败则降级为提取内嵌图片
    rendered = _render_pptx_to_screenshots(ppt_path, output_dir)
    rendered_map = {idx: path for idx, path in rendered}  # {slide_index: image_path}

    for slide_idx, slide in enumerate(prs.slides):
        slide_title = ""
        slide_texts = []
        slide_image_path = rendered_map.get(slide_idx + 1)  # 优先用全页截图

        for shape in slide.shapes:
            # 提取文字
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        is_title = False
                        try:
                            ph = shape.placeholder_format
                            if ph is not None and ph.idx == 0:
                                is_title = True
                        except Exception:
                            pass
                        if not slide_title and is_title:
                            slide_title = line
                        slide_texts.append(line)

            # 降级模式：若无全页截图，提取内嵌图片（每页只取第一张有意义的图）
            if not slide_image_path and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    img_bytes = image.blob
                    img = Image.open(io.BytesIO(img_bytes))
                    if img.width > 200 and img.height > 200:
                        img_filename = f"slide_{slide_idx+1:02d}.png"
                        img_path = os.path.join(output_dir, img_filename)
                        img.save(img_path, "PNG")
                        slide_image_path = img_path
                except Exception:
                    pass

        if slide_image_path and slide_image_path not in image_paths:
            image_paths.append(slide_image_path)

        slide_text = "\n".join(slide_texts)
        if slide_text:
            all_text_parts.append(f"【第{slide_idx+1}页】\n{slide_text}")

        slides_data.append({
            "index": slide_idx + 1,
            "title": slide_title or f"第{slide_idx+1}页",
            "text": slide_text,
            "has_image": slide_image_path is not None,
            "image_path": slide_image_path,
        })

    return {
        "slides": slides_data,
        "full_text": "\n\n".join(all_text_parts),
        "image_paths": image_paths,
    }


def parse_pdf(pdf_path: str) -> dict:
    """
    解析 PDF 文件（通常为 PPT 导出的 PDF），每页渲染为整页截图。
    返回与 parse_ppt 完全相同的结构。
    """
    import fitz  # PyMuPDF
    output_dir = os.path.join(os.path.dirname(pdf_path), "_ppt_images")
    os.makedirs(output_dir, exist_ok=True)

    doc = fitz.open(pdf_path)
    slides_data = []
    all_text_parts = []
    image_paths = []

    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        title = lines[0] if lines else f"第{i+1}页"

        img_filename = f"slide_{i+1:02d}.png"
        img_path = os.path.join(output_dir, img_filename)
        pix = page.get_pixmap(dpi=150)
        pix.save(img_path)
        image_paths.append(img_path)

        if text:
            all_text_parts.append(f"【第{i+1}页】\n{text}")

        slides_data.append({
            "index": i + 1,
            "title": title,
            "text": text,
            "has_image": True,
            "image_path": img_path,
        })

    doc.close()
    return {
        "slides": slides_data,
        "full_text": "\n\n".join(all_text_parts),
        "image_paths": image_paths,
    }


def parse_file(path: str) -> dict:
    """统一入口：按后缀路由到 parse_pdf 或 parse_ppt"""
    if path.lower().endswith(".pdf"):
        return parse_pdf(path)
    return parse_ppt(path)


def get_key_images(ppt_data: dict, max_images: int = 5) -> list:
    """
    从PPT中筛选关键图片（有图的页面，取前max_images张）
    """
    images = []
    for slide in ppt_data["slides"]:
        if slide["has_image"] and slide["image_path"]:
            images.append({
                "path": slide["image_path"],
                "slide_index": slide["index"],
                "slide_title": slide["title"],
            })
        if len(images) >= max_images:
            break
    return images
